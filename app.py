# Build marker 2026-06-15e — adds public Google Slides link import. Forces a fresh commit/redeploy.
"""
PPT -> ArcGIS StoryMaps Briefing converter (backend) — stateless build for Render.

Authentication: ArcGIS username + password
-------------------------------------------
The user signs in with their ArcGIS Online username and password. We validate the
credentials by signing in once, then keep them in an ENCRYPTED, http-only cookie
(Fernet, keyed off APP_SECRET_KEY) for the session — never in plain text, never on
the server. Each conversion re-authenticates from those credentials.

NOTE: username/password works only for plain ArcGIS accounts. Accounts that use
SSO/SAML or enforce multi-factor authentication cannot sign in this way (that is
what OAuth is for).

Design notes (unchanged from the OAuth build)
---------------------------------------------
* Stateless: no server-side session/job store. Any instance can serve any
  request; restarts/redeploys don't sign people out.
* Single-request streaming: the conversion runs inside one POST and streams
  newline-delimited JSON (NDJSON) progress to the browser. No background threads.
* Memory-aware: `arcgis` is heavy and Render Starter has 512 MB, so it's imported
  lazily and conversions are serialized with a semaphore.

Environment variables
---------------------
    APP_SECRET_KEY    (required in prod; derives the cookie-encryption key — keep stable)
    ARCGIS_PORTAL     (default https://www.arcgis.com; used as the default portal URL)
    SESSION_HOURS     (optional, default 4; how long a sign-in stays valid)

Why a Python backend at all: Esri publishes no REST API / data model for
StoryMaps Briefings; the only supported way to build one programmatically is the
ArcGIS API for Python (arcgis.apps.storymap.Briefing), which runs server-side.
"""

import io
import os
import json
import time
import base64
import hashlib
import secrets
import shutil
import tempfile
import threading
import traceback
import subprocess
import re
import urllib.parse

import requests

from cryptography.fernet import Fernet, InvalidToken
from flask import (
    Flask,
    request,
    jsonify,
    send_from_directory,
    Response,
    make_response,
    stream_with_context,
)

app = Flask(__name__, static_folder="static", static_url_path="")
app.secret_key = os.environ.get("APP_SECRET_KEY") or secrets.token_hex(32)

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

PORTAL_DEFAULT = os.environ.get("ARCGIS_PORTAL", "https://www.arcgis.com").rstrip("/")
SESSION_TTL = float(os.environ.get("SESSION_HOURS", "4")) * 3600

AUTH_COOKIE = "p2b_auth"

# Only one heavy (arcgis) build at a time per instance, to stay within 512 MB.
_build_gate = threading.BoundedSemaphore(1)


# ---------------------------------------------------------------------------
# Encrypted cookie session (stateless) — carries the ArcGIS credentials
# ---------------------------------------------------------------------------

def _fernet():
    """Derive a stable Fernet key from APP_SECRET_KEY."""
    digest = hashlib.sha256(app.secret_key.encode()).digest()
    return Fernet(base64.urlsafe_b64encode(digest))


def _horizon(bundle):
    return bundle.get("signed_in_at", 0) + SESSION_TTL


def read_auth():
    """Decrypt the auth cookie into a credentials bundle, or return None."""
    raw = request.cookies.get(AUTH_COOKIE)
    if not raw:
        return None
    try:
        bundle = json.loads(_fernet().decrypt(raw.encode()).decode())
    except (InvalidToken, ValueError):
        return None
    if _horizon(bundle) < time.time():
        return None
    return bundle


def write_auth(resp, bundle):
    token = _fernet().encrypt(json.dumps(bundle).encode()).decode()
    max_age = max(int(_horizon(bundle) - time.time()), 60)
    resp.set_cookie(
        AUTH_COOKIE, token,
        max_age=max_age, httponly=True, samesite="Lax",
        secure=request.is_secure,
    )


def clear_auth(resp):
    resp.delete_cookie(AUTH_COOKIE)


# ---------------------------------------------------------------------------
# ArcGIS sign-in (lazy import: arcgis is heavy)
# ---------------------------------------------------------------------------

def login_gis(portal, username, password):
    """
    Authenticate to ArcGIS and return (gis, display_name).
    Raises ModuleNotFoundError if the minimal install is missing a dependency,
    or a generic Exception if the credentials/portal are invalid.
    """
    from arcgis.gis import GIS
    gis = GIS(portal, username, password)
    me = gis.users.me
    display = getattr(me, "fullName", None) or getattr(me, "username", username)
    return gis, display


def get_gis(bundle):
    """Re-authenticate for a conversion from the stored credentials."""
    gis, _ = login_gis(bundle["portal"], bundle["username"], bundle["password"])
    return gis


# ---------------------------------------------------------------------------
# Routes: static page + config probe
# ---------------------------------------------------------------------------

@app.route("/")
def index():
    return send_from_directory(app.static_folder, "index.html")


@app.route("/api/config")
def config():
    return jsonify(portal=PORTAL_DEFAULT, auth="password")


@app.route("/healthz")
def healthz():
    return "ok", 200


# ---------------------------------------------------------------------------
# Routes: authentication
# ---------------------------------------------------------------------------

@app.route("/api/auth", methods=["POST"])
def auth():
    data = request.get_json(force=True, silent=True) or {}
    portal = (data.get("portal_url") or PORTAL_DEFAULT).strip().rstrip("/")
    username = (data.get("username") or "").strip()
    password = data.get("password") or ""

    if not username or not password:
        return jsonify(ok=False, error="Enter both a username and a password."), 400

    try:
        _, display = login_gis(portal, username, password)
    except ModuleNotFoundError as e:
        return jsonify(ok=False, error=f"Missing Python dependency '{e.name}'. The minimal "
                                       f"ArcGIS install doesn't include it. Add '{e.name}' to "
                                       f"requirements.txt and redeploy (see DEPLOY.md)."), 500
    except Exception as e:
        return jsonify(ok=False, error="Could not sign in. Check the portal URL, username, and "
                                       "password. Accounts that use SSO or multi-factor sign-in "
                                       f"can't use username/password. ({type(e).__name__})"), 401
    finally:
        password = None  # drop the plaintext reference promptly after this scope

    bundle = {
        "username": username,
        "password": data.get("password") or "",   # stored encrypted in the cookie below
        "portal": portal,
        "signed_in_at": time.time(),
    }
    resp = make_response(jsonify(ok=True, user=display, username=username,
                                 expires_at=_horizon(bundle)))
    write_auth(resp, bundle)
    return resp


@app.route("/api/session")
def session_status():
    bundle = read_auth()
    if not bundle:
        return jsonify(ok=True, signed_in=False)
    return jsonify(ok=True, signed_in=True, username=bundle["username"],
                   portal=bundle.get("portal", PORTAL_DEFAULT),
                   expires_at=_horizon(bundle), seconds_left=int(_horizon(bundle) - time.time()))


@app.route("/api/logout", methods=["POST"])
def logout():
    resp = make_response(jsonify(ok=True))
    clear_auth(resp)
    return resp


# ---------------------------------------------------------------------------
# PPTX parsing
# ---------------------------------------------------------------------------

def parse_pptx(file_bytes, workdir):
    """Extract per-slide title, body text, speaker notes, images, and SmartArt."""
    from pptx import Presentation
    from pptx.oxml.ns import qn

    # SmartArt (a "diagram") lives in a graphicFrame whose graphicData uri is this.
    DIAGRAM_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"

    prs = Presentation(io.BytesIO(file_bytes))
    slides = []

    for s_idx, slide in enumerate(prs.slides):
        title = ""
        body = []
        images = []
        smartart = []

        try:
            if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text_frame.text.strip()
        except Exception:
            pass

        for shape in slide.shapes:
            if shape.has_text_frame:
                if slide.shapes.title is not None and shape == slide.shapes.title:
                    continue
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if not txt and para.text:
                        txt = para.text.strip()
                    if txt:
                        body.append(txt)

            if shape.shape_type == 13 or getattr(shape, "image", None) is not None:
                try:
                    image = shape.image
                    ext = (image.ext or "png").lower()
                    fpath = os.path.join(workdir, f"slide{s_idx + 1}_{secrets.token_hex(4)}.{ext}")
                    with open(fpath, "wb") as fh:
                        fh.write(image.blob)
                    images.append({"path": fpath, "ext": ext})
                except Exception:
                    continue

            # SmartArt: a graphicFrame referencing the diagram namespace. We can't
            # read its content in Python, so we record its bounding box to render
            # later (see render_smartart_stream).
            try:
                gd = shape._element.find(".//" + qn("a:graphicData"))
                if gd is not None and DIAGRAM_URI in (gd.get("uri") or ""):
                    smartart.append({
                        "left": int(shape.left), "top": int(shape.top),
                        "width": int(shape.width), "height": int(shape.height),
                    })
            except Exception:
                pass

        notes = ""
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            pass

        if not title and body:
            title = body[0][:120]

        slides.append({
            "index": s_idx,
            "title": title or f"Slide {s_idx + 1}",
            "body": body,
            "notes": notes,
            "images": images,
            "smartart": smartart,
        })

    return slides


# ---------------------------------------------------------------------------
# Briefing construction (version-sensitive: the data model is undocumented)
# ---------------------------------------------------------------------------

def _log(message, level="info"):
    return {"type": "log", "level": level, "message": message}


# --- SmartArt rendering ----------------------------------------------------
# SmartArt can't be read as text/vector in Python, so we render it to an image
# with LibreOffice (pptx -> PDF) and crop the SmartArt's bounding box with
# PyMuPDF. This requires LibreOffice + pymupdf in the environment (the Docker
# build). When they're absent (e.g. the lean Render Starter native build), we
# skip SmartArt gracefully so everything else still works.

_EMU_PER_POINT = 12700
_SMARTART_DPI = 200


def _renderer_exe():
    """Return the LibreOffice executable if SmartArt rendering is possible, else None."""
    exe = shutil.which("soffice") or shutil.which("libreoffice")
    if not exe:
        return None
    try:
        import fitz  # noqa: F401  (PyMuPDF)
    except Exception:
        return None
    return exe


def render_smartart_stream(file_bytes, slides, workdir):
    """Render each SmartArt graphic to a PNG and append it to that slide's images.
    Generator: yields log events and mutates `slides` in place."""
    total = sum(len(s.get("smartart") or []) for s in slides)
    if total == 0:
        return

    exe = _renderer_exe()
    if not exe:
        yield _log(f"Found {total} SmartArt graphic(s), but this server can't render them "
                   f"(needs the LibreOffice-enabled build — see DEPLOY.md). Skipping.", "warn")
        return

    import fitz
    yield _log(f"Rendering {total} SmartArt graphic(s)…")

    pptx_path = os.path.join(workdir, "deck.pptx")
    with open(pptx_path, "wb") as fh:
        fh.write(file_bytes)

    # LibreOffice needs a writable HOME for its profile.
    env = dict(os.environ, HOME=workdir)
    try:
        subprocess.run([exe, "--headless", "--convert-to", "pdf", "--outdir", workdir, pptx_path],
                       env=env, capture_output=True, timeout=180, check=True)
    except Exception as e:
        yield _log(f"SmartArt rendering failed during PDF conversion ({type(e).__name__}); "
                   f"continuing without it.", "warn")
        return

    pdf_path = os.path.join(workdir, "deck.pdf")
    if not os.path.exists(pdf_path):
        yield _log("SmartArt rendering produced no PDF; continuing without it.", "warn")
        return

    zoom = fitz.Matrix(_SMARTART_DPI / 72.0, _SMARTART_DPI / 72.0)
    try:
        doc = fitz.open(pdf_path)
        for s in slides:
            sa = s.get("smartart") or []
            idx = s["index"]
            if not sa or idx >= doc.page_count:
                continue
            page = doc[idx]
            for j, bb in enumerate(sa):
                try:
                    clip = fitz.Rect(
                        bb["left"] / _EMU_PER_POINT, bb["top"] / _EMU_PER_POINT,
                        (bb["left"] + bb["width"]) / _EMU_PER_POINT,
                        (bb["top"] + bb["height"]) / _EMU_PER_POINT,
                    )
                    out = os.path.join(workdir, f"smartart_{idx + 1}_{j}_{secrets.token_hex(3)}.png")
                    page.get_pixmap(clip=clip, matrix=zoom).save(out)
                    s["images"].append({"path": out, "ext": "png"})
                    yield _log(f"  • rendered SmartArt on slide {idx + 1}")
                except Exception as e:
                    yield _log(f"  • SmartArt on slide {idx + 1} skipped ({type(e).__name__})", "warn")
        doc.close()
    except Exception as e:
        yield _log(f"SmartArt rendering error ({type(e).__name__}); continuing.", "warn")


# Flexible briefing slides accept 1–6 content blocks. We reserve block 0 for the
# slide's text and give each image its own block (a block holds only one image).
_MAX_IMAGE_BLOCKS = 5


def build_briefing_stream(gis, slides, briefing_title):
    """Generator: yields progress events while building the briefing.

    Briefing slide API (arcgis 2.4/2.5): create a flexible BriefingSlide with a
    given number of blocks, attach it with Briefing.add(slide=...), then fill
    each block via block.add_content(...). A block can hold multiple Text
    contents but only one Image, so each image gets its own block.
    """
    from arcgis.apps.storymap import Briefing
    from arcgis.apps.storymap import story_content as sc
    Text, Image, BriefingSlide = sc.Text, sc.Image, sc.BriefingSlide

    yield _log("Creating a new briefing draft in your ArcGIS Online content…")
    briefing = Briefing(gis=gis)

    cover_title = briefing_title or (slides[0]["title"] if slides else "Imported Briefing")

    for ps in slides:
        yield _log(f"Building slide {ps['index'] + 1}: “{ps['title']}”")

        images = ps["images"][:_MAX_IMAGE_BLOCKS]
        num_blocks = 1 + len(images)          # block 0 = text; one block per image

        slide = BriefingSlide(
            layout="flexible", num_blocks=num_blocks,
            title=ps["title"], story=briefing,
        )
        slide = briefing.add(slide=slide)     # attaches the slide and returns it
        blocks = slide.blocks

        # Block 0: body paragraphs and notes (a block may hold multiple Texts).
        if blocks:
            for para in ps["body"]:
                try:
                    blocks[0].add_content(Text(text=para))
                except Exception:
                    continue
            if ps["notes"]:
                try:
                    blocks[0].add_content(Text(text=f"Notes: {ps['notes']}"))
                except Exception:
                    pass

        # Remaining blocks: one image each.
        for i, img in enumerate(images, start=1):
            if i >= len(blocks):
                break
            try:
                blocks[i].add_content(Image(img["path"]))
                yield _log(f"  • uploaded image {os.path.basename(img['path'])}")
            except Exception as e:
                yield _log(f"  • image skipped ({type(e).__name__})", "warn")

    yield _log("Saving the briefing…")
    item = briefing.save(title=cover_title, publish=False)

    item_id = getattr(item, "id", None) or getattr(briefing, "_itemid", None)
    portal = gis.url.rstrip("/")
    yield {
        "type": "end",
        "status": "done",
        "result": {
            "item_id": item_id,
            "title": cover_title,
            "slide_count": len(slides),
            "item_url": f"{portal}/home/item.html?id={item_id}" if item_id else None,
            "edit_url": f"https://storymaps.arcgis.com/briefings/{item_id}/edit" if item_id else None,
        },
    }


def convert_stream(bundle, file_bytes, briefing_title):
    """Top-level generator for one conversion. Yields progress + final event."""
    workdir = tempfile.mkdtemp(prefix="p2b_")
    try:
        yield _log("Reading the PowerPoint file…")
        slides = parse_pptx(file_bytes, workdir)
        if not slides:
            yield {"type": "end", "status": "error", "error": "No slides were found in the file."}
            return
        yield _log(f"Found {len(slides)} slide(s).")

        # Serialize heavy arcgis builds to protect 512 MB instances.
        if not _build_gate.acquire(timeout=0.5):
            yield _log("Another conversion is in progress; waiting for it to finish…", "warn")
            _build_gate.acquire()
        try:
            yield from render_smartart_stream(file_bytes, slides, workdir)
            yield _log("Signing in to ArcGIS…")
            gis = get_gis(bundle)
            yield from build_briefing_stream(gis, slides, briefing_title)
        finally:
            _build_gate.release()
    except ModuleNotFoundError as e:
        traceback.print_exc()
        yield {"type": "end", "status": "error",
               "error": f"Missing Python dependency '{e.name}'. The minimal ArcGIS install "
                        f"doesn't include it. Add '{e.name}' to requirements.txt and redeploy "
                        f"(see DEPLOY.md → 'Minimal install')."}
    except Exception as e:
        traceback.print_exc()
        yield {"type": "end", "status": "error", "error": f"{type(e).__name__}: {e}"}
    finally:
        try:
            for f in os.listdir(workdir):
                os.remove(os.path.join(workdir, f))
            os.rmdir(workdir)
        except Exception:
            pass


def fetch_google_slides_pptx(url):
    """Download a PUBLIC Google Slides deck as .pptx bytes.

    Only docs.google.com links are accepted, and the export URL is built here
    (not taken from the user), which contains SSRF. Raises ValueError with a
    user-facing message on any problem.
    """
    parsed = urllib.parse.urlparse(url.strip())
    if parsed.scheme not in ("http", "https") or parsed.netloc not in (
        "docs.google.com", "www.docs.google.com"
    ):
        raise ValueError("Enter a Google Slides link from docs.google.com.")

    m = re.search(r"/presentation/d/(e/)?([a-zA-Z0-9_-]+)", parsed.path)
    if not m:
        raise ValueError("That doesn't look like a Google Slides link.")
    published, file_id = bool(m.group(1)), m.group(2)
    base = f"https://docs.google.com/presentation/d/{'e/' if published else ''}{file_id}"
    export_url = f"{base}/export/pptx"

    try:
        r = requests.get(export_url, timeout=60, allow_redirects=True)
    except Exception as e:
        raise ValueError(f"Could not reach Google Slides ({type(e).__name__}). Try again.")

    data = r.content or b""
    # A real .pptx is a ZIP (starts with 'PK\x03\x04'). Anything else (e.g. a
    # sign-in/permission HTML page) means the deck isn't publicly viewable.
    if r.status_code != 200 or data[:4] != b"PK\x03\x04":
        raise ValueError(
            "Couldn't download that deck. In Google Slides choose Share → General "
            "access → 'Anyone with the link' (Viewer), then paste the link again."
        )
    if len(data) > 100 * 1024 * 1024:
        raise ValueError("That deck is too large to convert.")
    return data


@app.route("/api/convert", methods=["POST"])
def convert():
    bundle = read_auth()
    if not bundle:
        return jsonify(ok=False, error="Your session has expired. Please sign in again."), 401

    upload = request.files.get("file")
    gslides_url = (request.form.get("gslides_url") or "").strip()

    if upload and upload.filename:
        if not upload.filename.lower().endswith(".pptx"):
            return jsonify(ok=False, error="Please upload a .pptx file. The older binary .ppt "
                                           "format isn't supported — Save As .pptx first."), 400
        file_bytes = upload.read()
    elif gslides_url:
        try:
            file_bytes = fetch_google_slides_pptx(gslides_url)
        except ValueError as e:
            return jsonify(ok=False, error=str(e)), 400
    else:
        return jsonify(ok=False, error="Upload a .pptx file or paste a Google Slides link."), 400

    briefing_title = (request.form.get("title") or "").strip()

    @stream_with_context
    def gen():
        for event in convert_stream(bundle, file_bytes, briefing_title):
            yield json.dumps(event) + "\n"

    resp = Response(gen(), mimetype="application/x-ndjson")
    resp.headers["Cache-Control"] = "no-cache"
    resp.headers["X-Accel-Buffering"] = "no"  # discourage proxy buffering
    return resp


if __name__ == "__main__":
    app.run(host="127.0.0.1", port=int(os.environ.get("PORT", 5000)), threaded=True)
